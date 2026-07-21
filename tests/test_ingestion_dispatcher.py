"""Tests for extension dispatch and the single-document ingestion CLI."""

from __future__ import annotations

import subprocess
import sys
from collections.abc import Callable
from email.message import EmailMessage
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from document_intelligence.ingestion.cli import main
from document_intelligence.ingestion.dispatcher import parse_document
from document_intelligence.ingestion.exceptions import (
    DocumentParseError,
    UnsupportedDocumentFormatError,
)
from document_intelligence.ingestion.models import BlockType, ParsedDocument, SourceFormat


def _write_pdf(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with path.open("wb") as handle:
        writer.write(handle)


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    )
    textbox.text = "Dispatcher presentation"
    presentation.save(path)


def _write_eml(path: Path, body: str = "Dispatcher email body") -> None:
    message = EmailMessage()
    message["Subject"] = "Dispatcher message"
    message["From"] = "sender@example.com"
    message["To"] = "recipient@example.com"
    message["Date"] = "Tue, 21 Jul 2026 10:00:00 +0000"
    message["Message-ID"] = "<dispatcher@example.com>"
    message.set_content(body)
    path.write_bytes(message.as_bytes())


def test_package_import_does_not_preload_cli_or_parsers() -> None:
    result = subprocess.run(
        [
            sys.executable,
            "-c",
            (
                "import sys\n"
                "import document_intelligence.ingestion\n"
                "names = ['cli', 'pdf_parser', 'pptx_parser', 'eml_parser']\n"
                "print(any(f'document_intelligence.ingestion.{n}' in sys.modules "
                "for n in names))\n"
            ),
        ],
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr
    assert result.stdout == "False\n"
    assert "RuntimeWarning" not in result.stderr


def test_package_level_exports_are_available() -> None:
    from document_intelligence.ingestion import (
        BlockType as exported_block_type,
        ParsedDocument as exported_document,
        parse_document as exported_parser,
    )

    assert exported_document is ParsedDocument
    assert exported_block_type is BlockType
    assert exported_parser is parse_document


@pytest.mark.parametrize(
    ("filename", "writer", "expected_format"),
    [
        ("fixture.pdf", _write_pdf, SourceFormat.PDF),
        ("fixture.pptx", _write_pptx, SourceFormat.PPTX),
        ("fixture.eml", _write_eml, SourceFormat.EML),
    ],
)
def test_dispatches_supported_formats(
    tmp_path: Path,
    filename: str,
    writer: Callable[[Path], None],
    expected_format: SourceFormat,
) -> None:
    path = tmp_path / filename
    writer(path)

    document = parse_document(path, "S-DISPATCH")

    assert document.source_format is expected_format


def test_dispatch_is_case_insensitive(tmp_path: Path) -> None:
    path = tmp_path / "UPPER.PDF"
    _write_pdf(path)

    assert parse_document(path).source_format is SourceFormat.PDF


@pytest.mark.parametrize("filename", ["unsupported.txt", "no_suffix"])
def test_dispatch_rejects_unsupported_or_missing_extension(
    tmp_path: Path,
    filename: str,
) -> None:
    path = tmp_path / filename
    path.write_text("unsupported", encoding="utf-8")

    with pytest.raises(UnsupportedDocumentFormatError):
        parse_document(path)


def test_dispatch_missing_supported_file_raises_parse_error(tmp_path: Path) -> None:
    with pytest.raises(DocumentParseError, match="input file does not exist"):
        parse_document(tmp_path / "missing.pdf")


def test_cli_writes_revalidatable_json_and_concise_summary(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    source_text = "FULL-DOCUMENT-TEXT-MUST-NOT-APPEAR-IN-SUMMARY"
    input_path = tmp_path / "input.eml"
    output_path = tmp_path / "nested" / "output.json"
    _write_eml(input_path, source_text)

    exit_code = main(
        [
            str(input_path),
            "--source-id",
            "S-CLI",
            "--output",
            str(output_path),
        ]
    )
    captured = capsys.readouterr()

    assert exit_code == 0
    assert captured.err == ""
    assert "source=S-CLI format=EML" in captured.out
    assert "blocks=" in captured.out
    assert "warnings=" in captured.out
    assert source_text not in captured.out
    document = ParsedDocument.model_validate_json(output_path.read_text("utf-8"))
    assert source_text in "\n".join(block.text for block in document.blocks)


def test_cli_prints_json_when_output_is_omitted(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    input_path = tmp_path / "stdout.eml"
    _write_eml(input_path)

    assert main([str(input_path), "--indent", "0"]) == 0
    captured = capsys.readouterr()

    ParsedDocument.model_validate_json(captured.out)
    assert captured.err == ""


def test_cli_returns_two_for_invalid_input_path(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    exit_code = main([str(tmp_path / "missing.eml")])
    captured = capsys.readouterr()

    assert exit_code == 2
    assert "input file does not exist" in captured.err


def test_module_cli_help_has_no_runpy_warning() -> None:
    result = subprocess.run(
        [
            sys.executable,
            "-m",
            "document_intelligence.ingestion.cli",
            "--help",
        ],
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr
    assert "usage:" in result.stdout
    assert "RuntimeWarning" not in result.stderr
    assert "found in sys.modules" not in result.stderr
    assert "prior to execution" not in result.stderr
