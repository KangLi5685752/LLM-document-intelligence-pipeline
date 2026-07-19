"""Tests for deterministic synthetic PPTX and EML challenge fixtures."""

from __future__ import annotations

import csv
import hashlib
import json
import re
import subprocess
import sys
import zipfile
from datetime import datetime
from email import policy
from email.parser import BytesParser
from email.utils import getaddresses
from pathlib import Path

import pytest
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
GENERATOR = ROOT / "scripts" / "generate_synthetic_corpus.py"
COMMITTED_ROOT = ROOT / "data" / "synthetic"
GROUND_TRUTH = ROOT / "data" / "manifests" / "synthetic_ground_truth.jsonl"
SOURCE_REGISTER = ROOT / "data" / "manifests" / "source_register.csv"
FOOTER = "Synthetic fixture – all names and details are fictional"
EMAIL_FOOTER = (
    "Synthetic fixture – all names, organisations, addresses, dates, budgets "
    "and events are fictional."
)
EXPECTED_RELATIVE_PATHS = {
    "pptx/S010_northstar_portfolio_update.pptx",
    "pptx/S011_civic_assist_steering_pack.pptx",
    "email/S012_atlas_migration_01.eml",
    "email/S013_atlas_migration_02.eml",
    "email/S014_atlas_migration_03.eml",
    "email/S015_harbour_analytics_01.eml",
    "email/S016_harbour_analytics_02.eml",
    "email/S017_harbour_analytics_03.eml",
}
PUBLIC_REGISTER_SNAPSHOT_SHA256 = (
    "90FC167CEE43E329D3E2D4011E1A3C44EEC03E91C09740596D945AC7F3DCDF73"
)


def _run_generator(output_root: Path, *, force: bool = False) -> subprocess.CompletedProcess[str]:
    command = [
        sys.executable,
        str(GENERATOR),
        "--output-root",
        str(output_root),
    ]
    if force:
        command.append("--force")
    return subprocess.run(
        command,
        cwd=ROOT,
        capture_output=True,
        text=True,
        check=False,
    )


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest().upper()


def _fixture_files(root: Path) -> list[Path]:
    return sorted(
        path
        for path in root.glob("*/*")
        if path.suffix.lower() in {".pptx", ".eml"}
    )


def _relative_paths(root: Path) -> set[str]:
    return {path.relative_to(root).as_posix() for path in _fixture_files(root)}


def _slide_text(slide: object) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            values.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def _parse_emails(root: Path) -> dict[str, object]:
    parser = BytesParser(policy=policy.default)
    return {
        path.stem.split("_", 1)[0]: parser.parsebytes(path.read_bytes())
        for path in sorted((root / "email").glob("*.eml"))
    }


@pytest.fixture(scope="module")
def generated_root(tmp_path_factory: pytest.TempPathFactory) -> Path:
    root = tmp_path_factory.mktemp("synthetic-corpus")
    result = _run_generator(root)
    assert result.returncode == 0, result.stderr
    return root


def test_generator_creates_exact_inventory_and_requires_force(
    generated_root: Path,
) -> None:
    assert _relative_paths(generated_root) == EXPECTED_RELATIVE_PATHS
    assert len(list((generated_root / "pptx").glob("*.pptx"))) == 2
    assert len(list((generated_root / "email").glob("*.eml"))) == 6

    refused = _run_generator(generated_root)
    assert refused.returncode != 0
    assert "rerun with --force" in refused.stderr

    replaced = _run_generator(generated_root, force=True)
    assert replaced.returncode == 0, replaced.stderr
    assert _relative_paths(generated_root) == EXPECTED_RELATIVE_PATHS


def test_generator_is_byte_deterministic(tmp_path: Path) -> None:
    first = tmp_path / "first"
    second = tmp_path / "second"
    first_result = _run_generator(first)
    second_result = _run_generator(second)
    assert first_result.returncode == second_result.returncode == 0

    first_paths = {
        path.relative_to(first).as_posix(): _sha256(path)
        for path in _fixture_files(first)
    }
    second_paths = {
        path.relative_to(second).as_posix(): _sha256(path)
        for path in _fixture_files(second)
    }
    assert first_paths == second_paths
    assert all(re.fullmatch(r"[0-9A-F]{64}", value) for value in first_paths.values())


def test_pptx_slide_counts_key_text_and_footers(generated_root: Path) -> None:
    s010 = Presentation(generated_root / "pptx/S010_northstar_portfolio_update.pptx")
    s011 = Presentation(generated_root / "pptx/S011_civic_assist_steering_pack.pptx")
    assert len(s010.slides) == 7
    assert len(s011.slides) == 8

    assert {
        "Programme status",
        "Green",
        "Emma Clarke",
        "GBP 1.35 million",
    } <= set(_slide_text(s010.slides[1]).splitlines())
    assert "Human review between Enrich and Publish" in _slide_text(s010.slides[3])
    assert "GBP 1.20 million" in _slide_text(s010.slides[5])
    assert "SUPERSEDED – do not use as the current programme position" in _slide_text(
        s010.slides[6]
    )

    assert "CivicAssist" in _slide_text(s011.slides[1])
    assert "Low confidence" in _slide_text(s011.slides[2])
    assert "No final launch date has been approved." in _slide_text(s011.slides[4])
    assert "GBP 880,000" in _slide_text(s011.slides[5])
    assert "Unresolved" in _slide_text(s011.slides[7])

    for presentation in (s010, s011):
        assert all(FOOTER in _slide_text(slide) for slide in presentation.slides)


def test_pptx_metadata_archive_and_relationships_are_controlled(
    generated_root: Path,
) -> None:
    expected = {
        "S010_northstar_portfolio_update.pptx": (
            "Northstar Programme Portfolio Update",
            datetime(2026, 7, 15),
        ),
        "S011_civic_assist_steering_pack.pptx": (
            "Civic Assist Pilot Steering Pack",
            datetime(2026, 7, 18),
        ),
    }
    for path in sorted((generated_root / "pptx").glob("*.pptx")):
        presentation = Presentation(path)
        title, document_date = expected[path.name]
        properties = presentation.core_properties
        assert properties.title == title
        assert properties.author == "Kang Li"
        assert properties.category == "Synthetic corpus fixture"
        assert properties.created == document_date
        assert properties.modified == document_date
        assert properties.last_modified_by == "Kang Li"
        assert all(
            not relationship.is_external
            for slide in presentation.slides
            for relationship in slide.part.rels.values()
        )

        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            assert [info.filename for info in infos] == sorted(
                info.filename for info in infos
            )
            assert all(info.date_time == (2026, 1, 1, 0, 0, 0) for info in infos)
            assert all(info.compress_type == zipfile.ZIP_DEFLATED for info in infos)
            assert not any(info.filename.startswith("ppt/media/") for info in infos)
            assert not any(
                b'TargetMode="External"' in archive.read(info.filename)
                for info in infos
                if info.filename.endswith(".rels")
            )


def test_eml_files_parse_with_unique_ids_example_addresses_and_footer(
    generated_root: Path,
) -> None:
    messages = _parse_emails(generated_root)
    assert set(messages) == {f"S{i:03d}" for i in range(12, 18)}
    message_ids = [message["Message-ID"] for message in messages.values()]
    assert len(message_ids) == len(set(message_ids)) == 6

    for message in messages.values():
        assert message.get_content_type() == "text/plain"
        body = message.get_content()
        assert EMAIL_FOOTER in body
        addresses = getaddresses(
            [
                str(message[header])
                for header in ("From", "To", "Cc")
                if message[header]
            ]
        )
        assert addresses
        assert all(address.lower().endswith("@example.com") for _, address in addresses)
        body_addresses = re.findall(r"[\w.+-]+@[\w.-]+", body)
        assert all(address.lower().endswith("@example.com") for address in body_addresses)


def test_eml_thread_headers_and_dates_are_exact(generated_root: Path) -> None:
    messages = _parse_emails(generated_root)
    assert messages["S012"]["Message-ID"] == "<atlas-001@example.com>"
    assert messages["S013"]["In-Reply-To"] == "<atlas-001@example.com>"
    assert messages["S013"]["References"] == "<atlas-001@example.com>"
    assert messages["S014"]["In-Reply-To"] == "<atlas-002@example.com>"
    assert messages["S014"]["References"] == (
        "<atlas-001@example.com> <atlas-002@example.com>"
    )
    assert messages["S016"]["In-Reply-To"] == "<harbour-001@example.com>"
    assert messages["S016"]["References"] == "<harbour-001@example.com>"
    assert messages["S017"]["In-Reply-To"] == "<harbour-002@example.com>"
    assert messages["S017"]["References"] == (
        "<harbour-001@example.com> <harbour-002@example.com>"
    )
    assert str(messages["S012"]["Date"]) == "Mon, 04 May 2026 09:00:00 +0000"
    assert str(messages["S017"]["Date"]) == "Thu, 04 Jun 2026 15:20:00 +0000"


def test_latest_emails_preserve_resolved_and_unresolved_history(
    generated_root: Path,
) -> None:
    messages = _parse_emails(generated_root)
    atlas = messages["S014"].get_content()
    assert "Status: Green" in atlas
    assert "Owner: Daniel Reed, effective 9 May 2026" in atlas
    assert "Approved budget: GBP 480,000" in atlas
    assert "Target migration date: 15 October 2026" in atlas
    assert "GBP 450,000 superseded" in atlas
    assert "30 September 2026" in atlas

    harbour = messages["S017"].get_content()
    assert "Status: Amber" in harbour
    assert "Approved budget: GBP 610,000" in harbour
    assert "12 November 2026 and 19 November 2026" in harbour
    assert "Final launch date: Unresolved; no final date is approved" in harbour
    assert "Vendor assurance risk: Open" in harbour


def test_ground_truth_has_valid_ids_splits_conflicts_and_evidence() -> None:
    records = [
        json.loads(line)
        for line in GROUND_TRUTH.read_text(encoding="utf-8").splitlines()
    ]
    assert len(records) == 4
    assert {record["challenge_id"] for record in records} == {
        f"GT-SYN-{index:03d}" for index in range(1, 5)
    }
    assert {
        record["document_family"]: record["provisional_split"] for record in records
    } == {
        "F-SYN-001": "development",
        "F-SYN-002": "held_out",
        "F-SYN-003": "development",
        "F-SYN-004": "held_out",
    }

    fact_ids: list[str] = []
    conflict_ids: list[str] = []
    group_ids: list[str] = []
    valid_sources = {f"S{index:03d}" for index in range(10, 18)}
    slide_counts = {"S010": 7, "S011": 8}
    committed_names = {path.stem.split("_", 1)[0] for path in _fixture_files(COMMITTED_ROOT)}

    for record in records:
        assert record["fictional"] is True
        assert set(record["source_ids"]) <= valid_sources
        current_ids = {
            item["fact_id"] for item in record["expected_current_facts"]
        }
        fact_ids.extend(current_ids)
        fact_ids.extend(item["fact_id"] for item in record["superseded_facts"])
        conflict_ids.extend(
            item["conflict_id"] for item in record["unresolved_conflicts"]
        )
        group_ids.extend(item["group_id"] for item in record["duplicate_groups"])

        for item in record["superseded_facts"]:
            assert item["superseded_by"] in current_ids
        for item in record["unresolved_conflicts"]:
            assert item["resolution"] is None
        for group in record["duplicate_groups"]:
            assert set(group["fact_ids"]) <= current_ids

        evidence_items = [
            evidence
            for key in (
                "expected_current_facts",
                "superseded_facts",
                "unresolved_conflicts",
            )
            for item in record[key]
            for evidence in item["evidence"]
        ]
        for evidence in evidence_items:
            source_id, location, *detail = evidence.split(":")
            assert source_id in record["source_ids"]
            assert source_id in committed_names
            if location == "slide":
                assert len(detail) == 1
                assert 1 <= int(detail[0]) <= slide_counts[source_id]
            else:
                assert location in {"body", "quoted_history"}
                assert not detail

    assert len(fact_ids) == len(set(fact_ids))
    assert len(conflict_ids) == len(set(conflict_ids))
    assert len(group_ids) == len(set(group_ids))


def test_source_register_preserves_public_rows_and_matches_fixtures() -> None:
    with SOURCE_REGISTER.open(encoding="utf-8", newline="") as handle:
        raw_rows = list(csv.reader(handle))
    header, data_rows = raw_rows[0], raw_rows[1:]
    rows = [dict(zip(header, row)) for row in data_rows]

    assert len(header) == 26
    assert all(len(row) == 26 for row in data_rows)
    assert [row["source_id"] for row in rows] == [
        f"S{index:03d}" for index in range(1, 18)
    ]
    canonical_public_rows = json.dumps(
        rows[:9],
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode()
    assert hashlib.sha256(canonical_public_rows).hexdigest().upper() == (
        PUBLIC_REGISTER_SNAPSHOT_SHA256
    )

    fixture_by_name = {path.name: path for path in _fixture_files(COMMITTED_ROOT)}
    for row in rows[9:]:
        assert row["local_filename"] in fixture_by_name
        assert row["sha256"] == _sha256(fixture_by_name[row["local_filename"]])
        assert row["original_file_committed"] == "true"
        assert row["redistribution_decision"] == "approved_for_redistribution"
        assert row["corpus_status"] == "approved"
        assert "fictional" in row["notes"].lower()
