"""Tests for slide- and shape-level PPTX ingestion."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from document_intelligence.ingestion.exceptions import DocumentParseError
from document_intelligence.ingestion.models import (
    BlockType,
    LocationType,
    ParsedDocument,
    SourceFormat,
)
from document_intelligence.ingestion.pptx_parser import parse_pptx


REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
S010_PATH = (
    REPOSITORY_ROOT
    / "data"
    / "synthetic"
    / "pptx"
    / "S010_northstar_portfolio_update.pptx"
)


def _new_presentation() -> Presentation:
    presentation = Presentation()
    while presentation.slides:
        relationship_id = presentation.slides._sldIdLst[-1].rId
        presentation.part.drop_rel(relationship_id)
        del presentation.slides._sldIdLst[-1]
    return presentation


def test_parse_pptx_extracts_title_and_one_based_slide_location(
    tmp_path: Path,
) -> None:
    path = tmp_path / "minimal.pptx"
    presentation = _new_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Temporary presentation"
    slide.placeholders[1].text = "Development fixture"
    presentation.save(path)

    document = parse_pptx(path, "S-TEMP-PPTX")

    assert document.source_format is SourceFormat.PPTX
    assert document.title == "Temporary presentation"
    assert document.metadata["slide_count"] == 1
    title_blocks = [
        block for block in document.blocks if block.block_type is BlockType.SLIDE_TITLE
    ]
    assert len(title_blocks) == 1
    assert title_blocks[0].text == "Temporary presentation"
    assert title_blocks[0].location.location_type is LocationType.SLIDE
    assert title_blocks[0].location.slide_number == 1
    ParsedDocument.model_validate_json(document.model_dump_json())


def test_parse_pptx_orders_shapes_by_position_then_original_index(
    tmp_path: Path,
) -> None:
    path = tmp_path / "ordering.pptx"
    presentation = _new_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    lower = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
    lower.text = "Lower shape"
    upper = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    upper.text = "Upper shape"
    presentation.save(path)

    document = parse_pptx(path, "S-ORDER")

    assert [block.text for block in document.blocks] == ["Upper shape", "Lower shape"]
    assert [block.metadata["original_shape_index"] for block in document.blocks] == [
        2,
        1,
    ]
    assert document.blocks[0].metadata["top"] < document.blocks[1].metadata["top"]
    assert all(block.location.element_index for block in document.blocks)


def test_parse_pptx_emits_one_deterministic_table_block(tmp_path: Path) -> None:
    path = tmp_path / "table.pptx"
    presentation = _new_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(2),
    ).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = ""
    table.cell(1, 0).text = ""
    table.cell(1, 1).text = "B"
    presentation.save(path)

    document = parse_pptx(path, "S-TABLE")

    table_blocks = [
        block for block in document.blocks if block.block_type is BlockType.TABLE
    ]
    assert len(table_blocks) == 1
    assert table_blocks[0].text == "A\t\n\tB"
    assert table_blocks[0].metadata["row_count"] == 2
    assert table_blocks[0].metadata["column_count"] == 2
    assert not any(
        block.block_type is BlockType.SHAPE_TEXT and block.text == table_blocks[0].text
        for block in document.blocks
    )


def test_parse_pptx_summarises_skipped_non_text_objects(tmp_path: Path) -> None:
    path = tmp_path / "visual.pptx"
    presentation = _new_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    presentation.save(path)

    document = parse_pptx(path, "S-VISUAL")

    assert document.metadata["skipped_object_count"] == 1
    assert len(document.warnings) == 1
    assert "Skipped 1" in document.warnings[0]


def test_parse_pptx_recurses_into_grouped_shapes(tmp_path: Path) -> None:
    path = tmp_path / "grouped.pptx"
    presentation = _new_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    child = group.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    child.text = "Grouped text"
    presentation.save(path)

    document = parse_pptx(path, "S-GROUP")
    grouped_blocks = [block for block in document.blocks if block.text == "Grouped text"]

    assert len(grouped_blocks) == 1
    assert grouped_blocks[0].metadata["group_path"]
    assert grouped_blocks[0].metadata["parent_shape_identifier"]


def test_parse_committed_s010_development_fixture() -> None:
    document = parse_pptx(S010_PATH, "S010")

    assert document.metadata["slide_count"] == 7
    assert any(
        block.location.slide_number == 2 and "Executive summary" in block.text
        for block in document.blocks
    )
    assert any(
        block.location.slide_number == 2 and "GBP 1.35 million" in block.text
        for block in document.blocks
    )
    assert all(block.location.slide_number >= 1 for block in document.blocks)
    ParsedDocument.model_validate_json(document.model_dump_json())


def test_parse_pptx_block_ids_are_deterministic_for_same_source() -> None:
    first = parse_pptx(S010_PATH, "S010")
    second = parse_pptx(S010_PATH, "S010")

    assert first.checksum_sha256 == second.checksum_sha256
    assert [block.block_id for block in first.blocks] == [
        block.block_id for block in second.blocks
    ]


@pytest.mark.parametrize("filename", ["missing.pptx", "malformed.pptx"])
def test_parse_pptx_missing_or_malformed_raises_structured_error(
    tmp_path: Path,
    filename: str,
) -> None:
    path = tmp_path / filename
    if filename == "malformed.pptx":
        path.write_bytes(b"not a presentation")

    with pytest.raises(DocumentParseError):
        parse_pptx(path)
