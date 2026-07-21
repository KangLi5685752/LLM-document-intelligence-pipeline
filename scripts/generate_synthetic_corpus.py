"""Generate deterministic fictional PPTX and EML challenge fixtures."""

from __future__ import annotations

import argparse
import hashlib
import sys
import zipfile
from datetime import datetime
from email import policy
from email.message import EmailMessage
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FOOTER = "Synthetic fixture – all names and details are fictional"
FICTIONAL_EMAIL_FOOTER = (
    "Synthetic fixture – all names, organisations, addresses, dates, budgets "
    "and events are fictional."
)
FIXED_ZIP_TIMESTAMP = (2026, 1, 1, 0, 0, 0)

EXPECTED_RELATIVE_PATHS = (
    Path("pptx/S010_northstar_portfolio_update.pptx"),
    Path("pptx/S011_civic_assist_steering_pack.pptx"),
    Path("email/S012_atlas_migration_01.eml"),
    Path("email/S013_atlas_migration_02.eml"),
    Path("email/S014_atlas_migration_03.eml"),
    Path("email/S015_harbour_analytics_01.eml"),
    Path("email/S016_harbour_analytics_02.eml"),
    Path("email/S017_harbour_analytics_03.eml"),
)

NAVY = RGBColor(27, 45, 71)
BLUE = RGBColor(45, 105, 155)
PALE_BLUE = RGBColor(224, 238, 247)
PALE_GREEN = RGBColor(224, 242, 232)
PALE_AMBER = RGBColor(255, 241, 204)
PALE_RED = RGBColor(249, 224, 224)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(35, 41, 48)
GREY = RGBColor(95, 103, 112)

S012_BODY = f"""Hello Atlas Delivery Team,

This is the initial fictional Atlas migration update.

- Status: Amber
- Owner: Priya Shah
- Proposed budget: GBP 450,000
- Target migration date: 30 September 2026
- Budget position: Proposed only; not approved

Regards,
Priya Shah
Atlas Migration Lead

{FICTIONAL_EMAIL_FOOTER}
"""

S013_BODY = f"""Hello Priya,

Finance has approved GBP 480,000 for the Atlas migration. The earlier
GBP 450,000 proposal is superseded. The target migration date changes to
15 October 2026. Status remains Amber, and Priya Shah remains the owner.

Regards,
Finance Review

{FICTIONAL_EMAIL_FOOTER}

----- Original Message -----
From: Priya Shah <priya.shah@example.com>
Sent: Mon, 04 May 2026 09:00:00 +0000
To: Atlas Delivery Team <atlas-team@example.com>
Subject: Atlas migration status
Message-ID: <atlas-001@example.com>

{S012_BODY}"""

S014_BODY = f"""Hello Atlas Delivery Team,

The current programme position is:

- Status: Green
- Owner: Daniel Reed, effective 9 May 2026
- Approved budget: GBP 480,000
- Target migration date: 15 October 2026

Priya Shah is the previous owner and is superseded by Daniel Reed. The
approved budget and target date from the finance reply remain current.

Regards,
Programme Office

{FICTIONAL_EMAIL_FOOTER}

Selected quoted history:
> Finance reply: GBP 480,000 approved; GBP 450,000 superseded.
> Finance reply: target date changed from 30 September 2026 to
> 15 October 2026.
> Initial update: status Amber; owner Priya Shah.
"""

S015_BODY = f"""Hello Harbour Steering Group,

This is the initial fictional Harbour analytics pilot update.

- Status: Green
- Owner: Maya Chen
- Working launch date: 12 November 2026
- Draft budget: GBP 620,000
- Budget position: Draft only; not approved

Regards,
Maya Chen
Harbour Analytics Lead

{FICTIONAL_EMAIL_FOOTER}
"""

S016_BODY = f"""Hello Maya,

The procurement ceiling is GBP 600,000. This ceiling is advisory and is
not a budget approval. The procurement dependency suggests 19 November
2026, while the delivery working date remains 12 November 2026. The date
conflict is unresolved.

Regards,
Procurement Desk

{FICTIONAL_EMAIL_FOOTER}

---------- Forwarded message ----------
From: Maya Chen <maya.chen@example.com>
Date: Mon, 01 Jun 2026 10:00:00 +0000
To: Harbour Steering Group <harbour-steering@example.com>
Subject: Harbour analytics pilot

> Status: Green
> Owner: Maya Chen
> Working launch date: 12 November 2026
> Draft budget: GBP 620,000; not approved
"""

S017_BODY = f"""Hello Harbour Steering Group,

The current steering position is:

- Status: Amber
- Owner: Maya Chen
- Approved budget: GBP 610,000
- Earlier draft budget GBP 620,000: Superseded
- Procurement ceiling GBP 600,000: Advisory, not an approval
- Candidate launch dates: 12 November 2026 and 19 November 2026
- Final launch date: Unresolved; no final date is approved
- Vendor assurance risk: Open

The latest message explicitly preserves the launch-date uncertainty.

Regards,
Steering Secretariat

{FICTIONAL_EMAIL_FOOTER}

Selected quoted and forwarded history:
> Initial working launch date: 12 November 2026.
> Procurement dependency date: 19 November 2026.
> Earlier draft budget: GBP 620,000.
> Procurement advisory ceiling: GBP 600,000.
"""


def calculate_sha256(path: Path) -> str:
    """Return an uppercase SHA-256 digest for a generated fixture."""
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest().upper()


def _set_text(
    shape: Any,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    **style: Any,
) -> Any:
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    _set_text(shape, text, **style)
    return shape


def _add_footer(slide: Any) -> None:
    _add_text(
        slide,
        FOOTER,
        0.45,
        7.08,
        12.4,
        0.2,
        font_size=9,
        color=GREY,
        alignment=PP_ALIGN.CENTER,
    )


def _new_slide(
    presentation: Any,
    title: str,
    *,
    title_size: int = 35,
) -> Any:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_text(
        slide,
        title,
        0.45,
        0.2,
        12.4,
        0.55,
        font_size=title_size,
        bold=True,
        color=NAVY,
    )
    _add_footer(slide)
    return slide


def _add_card(
    slide: Any,
    label: str,
    value: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = PALE_BLUE,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BLUE
    _set_text(
        shape,
        f"{label}\n{value}",
        font_size=16,
        bold=True,
        color=NAVY,
        alignment=PP_ALIGN.CENTER,
    )


def _add_table(
    slide: Any,
    headers: tuple[str, ...],
    rows: tuple[tuple[str, ...], ...],
    *,
    left: float = 0.55,
    top: float = 1.05,
    width: float = 12.2,
    height: float = 5.65,
) -> None:
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    ).table
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = WHITE
    for row_index, values in enumerate(rows, start=1):
        for column, value in enumerate(values):
            cell = table.cell(row_index, column)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else PALE_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(16)
                    run.font.color.rgb = DARK


def _add_process(
    slide: Any,
    labels: tuple[str, ...],
    *,
    top: float,
    start_left: float = 0.6,
    box_width: float = 2.1,
    gap: float = 0.45,
) -> list[Any]:
    positions = [
        start_left + index * (box_width + gap)
        for index in range(len(labels))
    ]
    for index in range(1, len(labels)):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(positions[index - 1] + box_width),
            Inches(top + 0.45),
            Inches(positions[index]),
            Inches(top + 0.45),
        )
        connector.line.color.rgb = BLUE
        connector.line.width = Pt(2)
        arrow_width = min(gap * 0.55, 0.3)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(positions[index] - gap / 2 - arrow_width / 2),
            Inches(top + 0.34),
            Inches(arrow_width),
            Inches(0.22),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BLUE
        arrow.line.color.rgb = BLUE

    boxes = []
    for left, label in zip(positions, labels):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(box_width),
            Inches(0.9),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PALE_BLUE
        box.line.color.rgb = BLUE
        _set_text(
            box,
            label,
            font_size=15,
            bold=True,
            color=NAVY,
            alignment=PP_ALIGN.CENTER,
        )
        boxes.append(box)
    return boxes


def _configure_presentation(
    presentation: Any,
    *,
    title: str,
    subject: str,
    document_date: datetime,
) -> None:
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    properties = presentation.core_properties
    properties.title = title
    properties.subject = subject
    properties.author = "Kang Li"
    properties.category = "Synthetic corpus fixture"
    properties.created = document_date
    properties.modified = document_date
    properties.last_modified_by = "Kang Li"


def _build_s010(path: Path) -> None:
    presentation = Presentation()
    _configure_presentation(
        presentation,
        title="Northstar Programme Portfolio Update",
        subject="Development synthetic challenge fixture",
        document_date=datetime(2026, 7, 15, 0, 0, 0),
    )

    slide = _new_slide(
        presentation,
        "Northstar Programme Portfolio Update",
        title_size=50,
    )
    _add_text(
        slide,
        "Reporting date: 15 July 2026",
        0.8,
        1.55,
        11.7,
        0.6,
        font_size=24,
        bold=True,
        color=BLUE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Development challenge fixture",
        0.8,
        2.35,
        11.7,
        0.5,
        font_size=24,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Fictional data created for portfolio testing; this is not a real record.",
        1.2,
        3.3,
        10.9,
        1.0,
        font_size=18,
        color=GREY,
        alignment=PP_ALIGN.CENTER,
    )

    slide = _new_slide(presentation, "Executive summary")
    _add_card(slide, "Programme status", "Green", 0.7, 1.2, 3.6, 1.25, fill=PALE_GREEN)
    _add_card(slide, "Programme owner", "Emma Clarke", 4.85, 1.2, 3.6, 1.25)
    _add_card(slide, "Approved budget", "GBP 1.35 million", 9.0, 1.2, 3.6, 1.25)
    _add_card(slide, "Target launch", "15 December 2026", 2.75, 3.25, 3.6, 1.25)
    _add_card(slide, "Reporting date", "15 July 2026", 6.95, 3.25, 3.6, 1.25)

    slide = _new_slide(presentation, "Workstream status table")
    _add_table(
        slide,
        ("Workstream", "Owner", "Status", "Target date"),
        (
            ("Data Platform", "Noah Bennett", "Amber", "30 September 2026"),
            ("Service Portal", "Chloe Hart", "Green", "15 November 2026"),
            ("Training and Adoption", "Aisha Cole", "Green", "1 December 2026"),
        ),
    )

    slide = _new_slide(presentation, "Delivery workflow")
    _add_text(
        slide,
        "Pilot workflow",
        0.65,
        1.1,
        3.0,
        0.45,
        font_size=18,
        bold=True,
        color=BLUE,
    )
    _add_process(slide, ("Ingest", "Validate", "Enrich", "Publish"), top=2.1)
    _add_text(
        slide,
        "Human review between Enrich and Publish",
        7.2,
        3.45,
        3.7,
        0.7,
        font_size=16,
        bold=True,
        color=NAVY,
        alignment=PP_ALIGN.CENTER,
    )
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(8.25),
        Inches(3.0),
        Inches(8.25),
        Inches(3.45),
    )
    connector.line.color.rgb = BLUE

    slide = _new_slide(presentation, "Risk register")
    _add_table(
        slide,
        ("Risk ID", "Risk", "Rating", "Owner", "Response"),
        (
            (
                "R-01",
                "Vendor integration delay",
                "Medium",
                "Noah Bennett",
                "Weekly integration checkpoint",
            ),
            (
                "R-02",
                "Source-data quality",
                "High",
                "Chloe Hart",
                "Validation rules and exception queue",
            ),
            (
                "R-03",
                "Training capacity",
                "Low",
                "Aisha Cole",
                "Train-the-trainer plan",
            ),
        ),
    )

    slide = _new_slide(presentation, "Approved decisions and change log")
    _add_table(
        slide,
        ("Change", "Previous value", "Current value", "Approval or effective date"),
        (
            (
                "Budget",
                "GBP 1.20 million",
                "GBP 1.35 million",
                "Approved 10 July 2026",
            ),
            (
                "Target launch",
                "30 November 2026",
                "15 December 2026",
                "Approved 12 July 2026",
            ),
            (
                "Programme owner",
                "Alex Morgan",
                "Emma Clarke",
                "Effective 14 July 2026",
            ),
        ),
    )

    slide = _new_slide(presentation, "Superseded snapshot – 30 June 2026")
    _add_text(
        slide,
        "SUPERSEDED – do not use as the current programme position",
        0.8,
        1.1,
        11.7,
        0.8,
        font_size=22,
        bold=True,
        color=RGBColor(153, 32, 32),
        alignment=PP_ALIGN.CENTER,
    )
    _add_card(slide, "Programme status", "Amber", 0.7, 2.45, 2.8, 1.35, fill=PALE_RED)
    _add_card(slide, "Programme owner", "Alex Morgan", 3.75, 2.45, 2.8, 1.35, fill=PALE_RED)
    _add_card(slide, "Budget", "GBP 1.20 million", 6.8, 2.45, 2.8, 1.35, fill=PALE_RED)
    _add_card(slide, "Target launch", "30 November 2026", 9.85, 2.45, 2.8, 1.35, fill=PALE_RED)

    _save_normalised_presentation(presentation, path)


def _build_s011(path: Path) -> None:
    presentation = Presentation()
    _configure_presentation(
        presentation,
        title="Civic Assist Pilot Steering Pack",
        subject="Held-out synthetic challenge fixture",
        document_date=datetime(2026, 7, 18, 0, 0, 0),
    )

    slide = _new_slide(
        presentation,
        "Civic Assist Pilot Steering Pack",
        title_size=50,
    )
    _add_text(
        slide,
        "Reporting date: 18 July 2026",
        0.8,
        1.55,
        11.7,
        0.6,
        font_size=24,
        bold=True,
        color=BLUE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Held-out challenge fixture",
        0.8,
        2.35,
        11.7,
        0.5,
        font_size=24,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Fictional data created for portfolio testing; this is not a real record.",
        1.2,
        3.3,
        10.9,
        1.0,
        font_size=18,
        color=GREY,
        alignment=PP_ALIGN.CENTER,
    )

    slide = _new_slide(presentation, "Scope and aliases")
    _add_card(slide, "Primary name", "Civic Assist", 0.7, 1.2, 3.7, 1.2)
    _add_card(slide, "Technical alias", "CivicAssist", 4.8, 1.2, 3.7, 1.2)
    _add_card(slide, "Project owner", "Lewis Grant", 8.9, 1.2, 3.7, 1.2)
    _add_card(slide, "Current status", "Amber", 0.7, 3.0, 3.7, 1.2, fill=PALE_AMBER)
    _add_text(
        slide,
        "Purpose: route fictional resident enquiries to appropriate knowledge "
        "and human-review paths",
        4.8,
        3.0,
        7.8,
        1.2,
        font_size=17,
        bold=True,
        color=NAVY,
        alignment=PP_ALIGN.CENTER,
    )

    slide = _new_slide(presentation, "Service workflow")
    _add_process(
        slide,
        ("Citizen query", "Triage", "Knowledge retrieval", "Human review", "Response"),
        top=1.75,
        start_left=0.35,
        box_width=2.1,
        gap=0.4,
    )
    exception_connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(6.4),
        Inches(2.65),
        Inches(6.4),
        Inches(4.0),
    )
    exception_connector.line.color.rgb = BLUE
    exception_connector.line.width = Pt(2)
    _add_process(
        slide,
        ("Low confidence", "Specialist queue"),
        top=4.0,
        start_left=5.0,
        box_width=2.5,
        gap=0.55,
    )
    _add_text(
        slide,
        "Exception path",
        6.55,
        3.25,
        2.1,
        0.4,
        font_size=14,
        bold=True,
        color=BLUE,
    )

    slide = _new_slide(presentation, "Pilot KPI table – synthetic figures")
    _add_table(
        slide,
        ("Measure", "Target", "Observed", "Interpretation"),
        (
            ("Routing accuracy", "85%", "82%", "Below target"),
            ("Human escalation rate", "<=15%", "14%", "Within threshold"),
            ("Median response time", "<=3.0 minutes", "2.8 minutes", "Within threshold"),
        ),
    )

    slide = _new_slide(presentation, "Schedule positions")
    _add_card(slide, "Operations working date", "31 March 2027", 1.0, 1.45, 4.8, 1.45, fill=PALE_AMBER)
    _add_card(slide, "Procurement dependency date", "15 April 2027", 7.5, 1.45, 4.8, 1.45, fill=PALE_AMBER)
    _add_text(
        slide,
        "No final launch date has been approved.",
        1.5,
        3.75,
        10.3,
        0.9,
        font_size=24,
        bold=True,
        color=RGBColor(153, 32, 32),
        alignment=PP_ALIGN.CENTER,
    )

    slide = _new_slide(presentation, "Budget decision")
    _add_card(slide, "Draft budget – superseded", "GBP 880,000", 1.0, 1.55, 4.7, 1.5, fill=PALE_RED)
    _add_card(slide, "Approved budget envelope", "GBP 920,000", 7.6, 1.55, 4.7, 1.5, fill=PALE_GREEN)
    _add_text(
        slide,
        "Approval date: 16 July 2026\nThe draft value is superseded.",
        2.0,
        3.8,
        9.3,
        1.0,
        font_size=20,
        bold=True,
        color=NAVY,
        alignment=PP_ALIGN.CENTER,
    )

    slide = _new_slide(presentation, "Risk matrix")
    _add_table(
        slide,
        ("Risk", "Rating", "State", "Owner"),
        (
            ("R-03 Vendor assurance", "High", "Open", "Lewis Grant"),
            ("R-04 Accessibility coverage", "Medium", "Mitigating", "Sofia Mills"),
            ("R-05 Knowledge freshness", "Medium", "Monitoring", "Ravi Dean"),
        ),
    )

    slide = _new_slide(presentation, "Steering summary")
    _add_card(slide, "Status", "Amber", 0.7, 1.2, 3.6, 1.15, fill=PALE_AMBER)
    _add_card(slide, "Project owner", "Lewis Grant", 4.85, 1.2, 3.6, 1.15)
    _add_card(slide, "Approved budget", "GBP 920,000", 9.0, 1.2, 3.6, 1.15)
    _add_card(slide, "Final launch date", "Unresolved", 0.7, 3.0, 3.6, 1.15, fill=PALE_RED)
    _add_card(slide, "Go/no-go review", "12 February 2027", 4.85, 3.0, 3.6, 1.15)
    _add_card(slide, "Vendor assurance risk", "Remains open", 9.0, 3.0, 3.6, 1.15, fill=PALE_RED)

    _save_normalised_presentation(presentation, path)


def _normalise_pptx_archive(source: Path, destination: Path) -> None:
    with zipfile.ZipFile(source, "r") as input_archive:
        members = [
            (name, input_archive.read(name))
            for name in sorted(input_archive.namelist())
        ]

    with zipfile.ZipFile(
        destination,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
    ) as output_archive:
        for name, content in members:
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIMESTAMP)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o600 << 16
            output_archive.writestr(
                info,
                content,
                compress_type=zipfile.ZIP_DEFLATED,
                compresslevel=9,
            )


def _save_normalised_presentation(presentation: Any, output_path: Path) -> None:
    raw_path = output_path.with_name(f".{output_path.name}.raw")
    normalised_path = output_path.with_name(f".{output_path.name}.normalised")
    try:
        presentation.save(raw_path)
        _normalise_pptx_archive(raw_path, normalised_path)
        normalised_path.replace(output_path)
    finally:
        raw_path.unlink(missing_ok=True)
        normalised_path.unlink(missing_ok=True)


def _write_email(
    output_path: Path,
    headers: tuple[tuple[str, str], ...],
    body: str,
) -> None:
    message = EmailMessage(policy=policy.SMTP)
    for name, value in headers:
        message[name] = value
    message.set_content(body, subtype="plain", charset="utf-8", cte="quoted-printable")

    temporary_path = output_path.with_name(f".{output_path.name}.tmp")
    try:
        temporary_path.write_bytes(message.as_bytes(policy=policy.SMTP))
        temporary_path.replace(output_path)
    finally:
        temporary_path.unlink(missing_ok=True)


def _email_definitions() -> tuple[tuple[str, tuple[tuple[str, str], ...], str], ...]:
    return (
        (
            "S012_atlas_migration_01.eml",
            (
                ("Date", "Mon, 04 May 2026 09:00:00 +0000"),
                ("From", "Priya Shah <priya.shah@example.com>"),
                ("To", "Atlas Delivery Team <atlas-team@example.com>"),
                ("Subject", "Atlas migration status"),
                ("Message-ID", "<atlas-001@example.com>"),
            ),
            S012_BODY,
        ),
        (
            "S013_atlas_migration_02.eml",
            (
                ("Date", "Wed, 06 May 2026 14:30:00 +0000"),
                ("From", "Finance Review <finance.review@example.com>"),
                ("To", "Priya Shah <priya.shah@example.com>"),
                ("Cc", "Atlas Delivery Team <atlas-team@example.com>"),
                ("Subject", "Atlas migration status"),
                ("Message-ID", "<atlas-002@example.com>"),
                ("In-Reply-To", "<atlas-001@example.com>"),
                ("References", "<atlas-001@example.com>"),
            ),
            S013_BODY,
        ),
        (
            "S014_atlas_migration_03.eml",
            (
                ("Date", "Fri, 08 May 2026 16:15:00 +0000"),
                ("From", "Programme Office <programme.office@example.com>"),
                ("To", "Atlas Delivery Team <atlas-team@example.com>"),
                ("Subject", "Atlas migration status"),
                ("Message-ID", "<atlas-003@example.com>"),
                ("In-Reply-To", "<atlas-002@example.com>"),
                ("References", "<atlas-001@example.com> <atlas-002@example.com>"),
            ),
            S014_BODY,
        ),
        (
            "S015_harbour_analytics_01.eml",
            (
                ("Date", "Mon, 01 Jun 2026 10:00:00 +0000"),
                ("From", "Maya Chen <maya.chen@example.com>"),
                ("To", "Harbour Steering Group <harbour-steering@example.com>"),
                ("Subject", "Harbour analytics pilot"),
                ("Message-ID", "<harbour-001@example.com>"),
            ),
            S015_BODY,
        ),
        (
            "S016_harbour_analytics_02.eml",
            (
                ("Date", "Tue, 02 Jun 2026 11:45:00 +0000"),
                ("From", "Procurement Desk <procurement@example.com>"),
                ("To", "Maya Chen <maya.chen@example.com>"),
                ("Subject", "Harbour analytics pilot"),
                ("Message-ID", "<harbour-002@example.com>"),
                ("In-Reply-To", "<harbour-001@example.com>"),
                ("References", "<harbour-001@example.com>"),
            ),
            S016_BODY,
        ),
        (
            "S017_harbour_analytics_03.eml",
            (
                ("Date", "Thu, 04 Jun 2026 15:20:00 +0000"),
                ("From", "Steering Secretariat <steering.secretariat@example.com>"),
                ("To", "Harbour Steering Group <harbour-steering@example.com>"),
                ("Subject", "Harbour analytics pilot"),
                ("Message-ID", "<harbour-003@example.com>"),
                ("In-Reply-To", "<harbour-002@example.com>"),
                ("References", "<harbour-001@example.com> <harbour-002@example.com>"),
            ),
            S017_BODY,
        ),
    )


def generate_synthetic_corpus(
    output_root: Path,
    *,
    force: bool = False,
) -> list[Path]:
    """Generate all expected deterministic synthetic fixtures."""
    output_root = output_root.resolve()
    expected_paths = [output_root / path for path in EXPECTED_RELATIVE_PATHS]
    existing = [path for path in expected_paths if path.exists()]
    if existing and not force:
        names = ", ".join(str(path.relative_to(output_root)) for path in existing)
        raise FileExistsError(
            f"expected output already exists: {names}; rerun with --force"
        )

    (output_root / "pptx").mkdir(parents=True, exist_ok=True)
    (output_root / "email").mkdir(parents=True, exist_ok=True)

    _build_s010(output_root / EXPECTED_RELATIVE_PATHS[0])
    _build_s011(output_root / EXPECTED_RELATIVE_PATHS[1])
    for filename, headers, body in _email_definitions():
        _write_email(output_root / "email" / filename, headers, body)

    return expected_paths


def _build_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Generate deterministic fictional PPTX and EML fixtures.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path("data/synthetic"),
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Replace existing expected synthetic outputs.",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    """Generate fixtures from command-line arguments."""
    args = _build_argument_parser().parse_args(argv)
    try:
        generated = generate_synthetic_corpus(args.output_root, force=args.force)
    except (OSError, ValueError, zipfile.BadZipFile) as error:
        print(f"error: {error}", file=sys.stderr)
        return 1

    root = args.output_root.resolve()
    print(f"generated={len(generated)} output_root={root}")
    for path in generated:
        print(
            f"{path.relative_to(root).as_posix()} "
            f"sha256={calculate_sha256(path)}"
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
