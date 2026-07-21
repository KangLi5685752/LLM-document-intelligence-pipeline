"""Slide- and shape-level PPTX ingestion with deterministic ordering."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation

from document_intelligence.ingestion.exceptions import DocumentParseError
from document_intelligence.ingestion.models import (
    BlockType,
    DocumentBlock,
    LocationType,
    ParseStatus,
    ParsedDocument,
    SourceFormat,
    SourceLocation,
)
from document_intelligence.ingestion.utils import (
    calculate_sha256,
    make_block_id,
    make_document_id,
    normalize_text,
    validate_input_file,
)


def _simple_property(value: Any) -> str | None:
    """Convert a core property into a stable serialisable string."""
    if value is None:
        return None
    if isinstance(value, datetime):
        return value.isoformat()
    normalized = normalize_text(str(value))
    return normalized or None


def _shape_number(shape: Any, attribute: str) -> int:
    """Return one geometry value as a plain integer."""
    value = getattr(shape, attribute, 0)
    return int(value or 0)


def _ordered_shapes(shapes: Any) -> list[tuple[int, Any]]:
    """Return 1-based original shape indexes in stable spatial order."""
    indexed = list(enumerate(shapes, start=1))
    return sorted(
        indexed,
        key=lambda item: (
            _shape_number(item[1], "top"),
            _shape_number(item[1], "left"),
            item[0],
        ),
    )


def _is_group_shape(shape: Any) -> bool:
    """Identify group shapes without depending on enum display strings."""
    return hasattr(shape, "shapes") and not getattr(shape, "has_table", False)


def _is_title_shape(shape: Any, title_shape_id: int | None) -> bool:
    return title_shape_id is not None and getattr(shape, "shape_id", None) == title_shape_id


def _table_text(shape: Any) -> tuple[str, int, int]:
    """Represent a table as tab-separated rows while preserving empty cells."""
    table = shape.table
    rows = [
        "\t".join(normalize_text(cell.text) for cell in row.cells)
        for row in table.rows
    ]
    return "\n".join(rows), len(table.rows), len(table.columns)


def parse_pptx(path: Path, source_id: str | None = None) -> ParsedDocument:
    """Parse one PPTX into ordered title, text-shape, and table blocks."""
    resolved_path = validate_input_file(Path(path))

    try:
        checksum = calculate_sha256(resolved_path)
        document_id = make_document_id(source_id, checksum)
        presentation = Presentation(resolved_path)
        blocks: list[DocumentBlock] = []
        skipped_types: dict[str, int] = {}
        fallback_title: str | None = None

        def emit_shape(
            shape: Any,
            *,
            slide_number: int,
            original_shape_index: int,
            title_shape_id: int | None,
            group_path: list[str],
            parent_shape_identifier: str | None,
        ) -> None:
            nonlocal fallback_title

            shape_name = str(getattr(shape, "name", "unnamed shape"))
            shape_identifier = f"slide-{slide_number}-shape-{getattr(shape, 'shape_id', original_shape_index)}"

            if _is_group_shape(shape):
                child_group_path = [*group_path, shape_identifier]
                for child_index, child in _ordered_shapes(shape.shapes):
                    emit_shape(
                        child,
                        slide_number=slide_number,
                        original_shape_index=child_index,
                        title_shape_id=title_shape_id,
                        group_path=child_group_path,
                        parent_shape_identifier=shape_identifier,
                    )
                return

            common_metadata: dict[
                str, str | int | float | bool | None | list[str]
            ] = {
                "shape_name": shape_name,
                "original_shape_index": original_shape_index,
                "top": _shape_number(shape, "top"),
                "left": _shape_number(shape, "left"),
                "width": _shape_number(shape, "width"),
                "height": _shape_number(shape, "height"),
                "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                "is_table": bool(getattr(shape, "has_table", False)),
                "shape_type": str(getattr(shape, "shape_type", "unknown")),
            }
            if group_path:
                common_metadata["group_path"] = group_path
            if parent_shape_identifier is not None:
                common_metadata["parent_shape_identifier"] = parent_shape_identifier

            block_type: BlockType | None = None
            text = ""
            if getattr(shape, "has_table", False):
                block_type = BlockType.TABLE
                text, row_count, column_count = _table_text(shape)
                common_metadata["row_count"] = row_count
                common_metadata["column_count"] = column_count
            elif getattr(shape, "has_text_frame", False):
                text = normalize_text(getattr(shape, "text", "") or "")
                if text:
                    if _is_title_shape(shape, title_shape_id):
                        block_type = BlockType.SLIDE_TITLE
                        if fallback_title is None:
                            fallback_title = text
                    else:
                        block_type = BlockType.SHAPE_TEXT

            if block_type is None:
                shape_type = str(getattr(shape, "shape_type", "unknown"))
                skipped_types[shape_type] = skipped_types.get(shape_type, 0) + 1
                return

            sequence = len(blocks) + 1
            blocks.append(
                DocumentBlock(
                    block_id=make_block_id(document_id, sequence),
                    sequence=sequence,
                    block_type=block_type,
                    text=text,
                    location=SourceLocation(
                        location_type=LocationType.SLIDE,
                        location_value=str(slide_number),
                        slide_number=slide_number,
                        element_index=original_shape_index,
                    ),
                    metadata=common_metadata,
                )
            )

        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            title_shape_id = (
                int(title_shape.shape_id) if title_shape is not None else None
            )
            for original_shape_index, shape in _ordered_shapes(slide.shapes):
                emit_shape(
                    shape,
                    slide_number=slide_number,
                    original_shape_index=original_shape_index,
                    title_shape_id=title_shape_id,
                    group_path=[],
                    parent_shape_identifier=None,
                )

        core = presentation.core_properties
        core_title = _simple_property(core.title)
        author = _simple_property(core.author)
        warnings: list[str] = []
        skipped_count = sum(skipped_types.values())
        if skipped_count:
            warnings.append(
                f"Skipped {skipped_count} unsupported or non-text slide objects"
            )

        metadata: dict[str, str | int | float | bool | None | list[str]] = {
            "slide_count": len(presentation.slides),
            "skipped_object_count": skipped_count,
            "skipped_object_types": [
                f"{shape_type}:{count}"
                for shape_type, count in sorted(skipped_types.items())
            ],
        }
        for key in (
            "title",
            "subject",
            "author",
            "category",
            "created",
            "modified",
            "last_modified_by",
        ):
            value = _simple_property(getattr(core, key, None))
            if value is not None:
                metadata[f"core_{key}"] = value

        document_date = core.created if isinstance(core.created, datetime) else None
        return ParsedDocument(
            document_id=document_id,
            source_id=source_id,
            source_format=SourceFormat.PPTX,
            filename=resolved_path.name,
            checksum_sha256=checksum,
            title=core_title or fallback_title,
            document_date=document_date,
            authors_or_senders=[author] if author else [],
            blocks=blocks,
            metadata=metadata,
            warnings=warnings,
            parse_status=(
                ParseStatus.SUCCESS_WITH_WARNINGS
                if warnings
                else ParseStatus.SUCCESS
            ),
        )
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError(
            "unable to parse PPTX",
            path=resolved_path,
            source_format=SourceFormat.PPTX,
        ) from error
